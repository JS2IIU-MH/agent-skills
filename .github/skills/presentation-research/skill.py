import os
import datetime
from typing import List, Dict
import re
import getpass

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None

import requests
from bs4 import BeautifulSoup


class PresentationResearchSkill:
    """Minimal Python implementation scaffold for the presentation-research skill.

    Notes:
    - This implementation uses a mocked web_search (placeholder).
    - Replace `mock_web_search` with the platform's `web_search` integration.
    - Uses `python-pptx` to generate a simple .pptx file.
    """

    def __init__(self, output_dir: str = "outputs"):
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def mock_web_search(self, query: str, top_n: int = 3) -> List[Dict]:
        # Retain mock as fallback
        return [
            {
                "title": f"Result for {query} - {i+1}",
                "url": f"https://example.com/search?q={query.replace(' ', '+')}&i={i+1}",
                "excerpt": f"Short excerpt about {query} (result {i+1}).",
            }
            for i in range(top_n)
        ]

    def web_search(self, query: str, top_n: int = 5, timeout: int = 10) -> List[Dict]:
        """Perform a simple web search using DuckDuckGo HTML endpoint and parse results.

        Falls back to `mock_web_search` on errors.
        """
        try:
            params = {"q": query}
            headers = {
                "User-Agent": "Mozilla/5.0 (compatible; PresentationResearchSkill/1.0)"
            }
            res = requests.post("https://html.duckduckgo.com/html/", data=params, headers=headers, timeout=timeout)
            res.raise_for_status()
            soup = BeautifulSoup(res.text, "html.parser")
            results = []
            for a in soup.select(".result__a")[:top_n]:
                title = a.get_text().strip()
                href = a.get('href')
                # DuckDuckGo returns redirect links; try to extract actual URL if present
                url = href or ""
                # Try to find snippet nearby
                parent = a.find_parent()
                excerpt = ""
                if parent:
                    snippet = parent.select_one('.result__snippet')
                    if snippet:
                        excerpt = snippet.get_text().strip()
                results.append({"title": title, "url": url, "excerpt": excerpt})
            if not results:
                return self.mock_web_search(query, top_n=top_n)
            return results
        except Exception:
            return self.mock_web_search(query, top_n=top_n)

    def generate_pptx(self, prompt: str, language: str = "ja", slide_count: int = 10, filename: str = None) -> str:
        if Presentation is None:
            raise RuntimeError("python-pptx is not installed. See requirements.txt")

        now = datetime.datetime.utcnow().strftime("%Y-%m-%d_%H%M%S")
        if not filename:
            safe_name = "presentation"
            filename = f"{safe_name}_{now}.pptx"

        out_path = os.path.join(self.output_dir, filename)

        # Try to load a template PPTX from templates/template.pptx next to this file.
        template_path = os.path.join(os.path.dirname(__file__), "templates", "template.pptx")
        if os.path.exists(template_path):
            try:
                prs = Presentation(template_path)
            except Exception:
                prs = Presentation()
        else:
            prs = Presentation()

        # If the template contains placeholders like {{TITLE}}, {{DATE}}, {{LANG}}, {{SUBTITLE}}, {{AUTHOR}}, replace them.
        def apply_template_placeholders(prs_obj, mapping: Dict[str, str]):
            pattern = re.compile(r"\{\{\s*([A-Z0-9_]+)\s*\}\}")
            for slide in prs_obj.slides:
                for shape in slide.shapes:
                    try:
                        if not getattr(shape, 'has_text_frame', False):
                            continue
                        tf = shape.text_frame
                    except Exception:
                        continue
                    for para in tf.paragraphs:
                        # replace in each run to preserve formatting when possible
                        for run in list(para.runs):
                            text = run.text
                            if not text:
                                continue
                            def _repl(m):
                                key = m.group(1).upper()
                                return mapping.get(key, m.group(0))
                            new_text = pattern.sub(_repl, text)
                            if new_text != text:
                                run.text = new_text

        now_human = datetime.datetime.now().strftime("%Y-%m-%d")
        mapping = {
            "TITLE": prompt if len(prompt) < 100 else prompt[:97] + "...",
            "DATE": now_human,
            "LANG": language,
            "SUBTITLE": f"Generated: {now} ({language})",
            "AUTHOR": getpass.getuser() if hasattr(getpass, 'getuser') else "",
        }
        try:
            apply_template_placeholders(prs, mapping)
        except Exception:
            pass

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        if title_shape is not None:
            title_shape.text = prompt if len(prompt) < 100 else prompt[:97] + "..."
        else:
            left = Inches(0.5)
            top = Inches(0.3)
            width = Inches(9)
            height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = prompt if len(prompt) < 100 else prompt[:97] + "..."
        # Some templates do not have placeholder idx 1 for subtitle; handle gracefully
        try:
            subtitle = slide.placeholders[1]
            subtitle.text = f"Generated: {now} ({language})"
        except Exception:
            left = Inches(1)
            top = Inches(1.7)
            width = Inches(8)
            height = Inches(0.8)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = f"Generated: {now} ({language})"

        # Helper to select a safe slide layout index
        def pick_layout(idx: int):
            try:
                return prs.slide_layouts[idx]
            except Exception:
                return prs.slide_layouts[0]

        def set_slide_title(slide, text: str):
            title_shape = slide.shapes.title
            if title_shape is not None:
                try:
                    title_shape.text = text
                    return
                except Exception:
                    pass
            # fallback: add textbox
            left = Inches(0.5)
            top = Inches(0.2)
            width = Inches(9)
            height = Inches(0.9)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = text

        def get_body_text_frame(slide):
            # Prefer placeholder index 1, else any placeholder with text_frame, else create textbox
            try:
                if len(slide.placeholders) > 1:
                    return slide.placeholders[1].text_frame
            except Exception:
                pass
            for ph in slide.placeholders:
                try:
                    if hasattr(ph, 'text_frame'):
                        return ph.text_frame
                except Exception:
                    continue
            # fallback create textbox
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(4.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            return txBox.text_frame

        # Table of contents
        toc = prs.slides.add_slide(pick_layout(1))
        set_slide_title(toc, "目次" if language.startswith("ja") else "Table of Contents")
        body = get_body_text_frame(toc)
        sections = [
            "エグゼクティブサマリー" if language.startswith("ja") else "Executive Summary",
            "調査結果" if language.startswith("ja") else "Research Findings",
            "詳細分析" if language.startswith("ja") else "Detailed Analysis",
            "結論・提言" if language.startswith("ja") else "Conclusions & Recommendations",
            "参考文献" if language.startswith("ja") else "References",
        ]
        for s in sections:
            p = body.add_paragraph()
            p.text = s

        # Mock research + content slides
        research = self.mock_web_search(prompt, top_n=5)

        # Executive summary
        s = prs.slides.add_slide(pick_layout(1))
        set_slide_title(s, sections[0])
        tf = get_body_text_frame(s)
        tf.text = "主要な発見:" if language.startswith("ja") else "Key findings:"
        for r in research[:3]:
            p = tf.add_paragraph()
            p.text = f"- {r['title']} ({r['url']})"

        # Research findings slides
        for i, r in enumerate(research):
            s = prs.slides.add_slide(pick_layout(1))
            set_slide_title(s, f"{sections[1]} - {i+1}")
            tf = get_body_text_frame(s)
            tf.text = r.get('excerpt', '')
            p = tf.add_paragraph()
            p.text = r.get('url', '')

        # Conclusion
        s = prs.slides.add_slide(pick_layout(1))
        set_slide_title(s, sections[3])
        tf = get_body_text_frame(s)
        tf.text = ("結論: 調査に基づく推奨事項をここに記載してください。" if language.startswith("ja")
               else "Conclusion: Add actionable recommendations based on the research.")

        # References
        s = prs.slides.add_slide(pick_layout(1))
        set_slide_title(s, sections[4])
        tf = get_body_text_frame(s)
        tf.text = "References:"
        for r in research:
            p = tf.add_paragraph()
            p.text = f"- {r.get('title', '')}: {r.get('url', '')}"

        prs.save(out_path)
        return out_path


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Run minimal presentation-research skill demo")
    parser.add_argument("prompt", nargs="?", default=None, help="Prompt / theme for the presentation")
    parser.add_argument("--lang", default="ja", help="Language code (ja/en)")
    parser.add_argument("--slides", type=int, default=10, help="Approximate number of slides")
    args = parser.parse_args()

    prompt = args.prompt or "生成AIの2025年のトレンドについてプレゼン資料を作ってください。"
    skill = PresentationResearchSkill(output_dir=os.path.join(os.getcwd(), "outputs"))
    out = skill.generate_pptx(prompt, language=args.lang, slide_count=args.slides)
    print(f"Generated: {out}")
