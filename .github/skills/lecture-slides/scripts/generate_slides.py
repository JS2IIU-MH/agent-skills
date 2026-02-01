import argparse
import json
import logging
from pptx import Presentation
from pptx.util import Inches

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def create_presentation(data_file, output_file, template_file=None):
    """
    Creates a PowerPoint presentation from a JSON data file.
    
    Args:
        data_file (str): Path to the JSON file containing slide data.
        output_file (str): Path to save the generated presentation.
        template_file (str, optional): Path to a PowerPoint template file.
    """
    try:
        # Load data
        with open(data_file, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # Initialize presentation
        if template_file:
            logging.info(f"Using template: {template_file}")
            prs = Presentation(template_file)
        else:
            logging.info("Creating new presentation without template")
            prs = Presentation()

        # Iterate through slides
        for i, slide_data in enumerate(data.get('slides', [])):
            # Determine layout
            # Default logic: First slide is Title Slide (0), others are Title and Content (1)
            # This works for default templates. Custom templates might vary.
            layout_index = slide_data.get('layout_index')
            
            if layout_index is None:
                if i == 0:
                    layout_index = 0 # Title Slide
                else:
                    layout_index = 1 # Title and Content
            
            # Check if layout index exists
            if layout_index >= len(prs.slide_layouts):
                logging.warning(f"Layout index {layout_index} out of range. Using 0.")
                layout_index = 0
                
            slide_layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set Title
            title = slide_data.get('title', '')
            if slide.shapes.title:
                slide.shapes.title.text = title
            
            # Set Content (Body)
            content = slide_data.get('content', [])
            # Find the body placeholder (usually index 1, but we can verify)
            # In standard layouts, placeholders[1] is usually the content body.
            if len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                if body_shape.has_text_frame:
                    tf = body_shape.text_frame
                    tf.clear() # Clear existing text (e.g. "Click to add text")
                    
                    if isinstance(content, list):
                        for point in content:
                            p = tf.add_paragraph()
                            p.text = str(point)
                            p.level = 0
                    elif isinstance(content, str):
                        p = tf.add_paragraph()
                        p.text = content
            
            # Set Speaker Notes
            notes = slide_data.get('notes', '')
            if notes:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = notes

        # Save presentation
        prs.save(output_file)
        logging.info(f"Presentation saved to: {output_file}")
        
    except Exception as e:
        logging.error(f"Failed to create presentation: {e}")
        raise

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate PowerPoint from JSON data.")
    parser.add_argument("data_file", help="Path to the JSON data file")
    parser.add_argument("output_file", help="Path to the output .pptx file")
    parser.add_argument("--template", help="Path to the .pptx template file", default=None)
    
    args = parser.parse_args()
    
    create_presentation(args.data_file, args.output_file, args.template)
